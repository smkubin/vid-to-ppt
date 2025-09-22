import os
import sys
import tempfile
import cv2
from datetime import datetime
from typing import List, Tuple

import numpy as np
from PIL import Image
from moviepy import VideoFileClip
from pptx import Presentation
from pptx.util import Inches
import whisper
from skimage.metrics import structural_similarity as ssim

# -----------------------------------------------------------------------------
# CONFIGURATION
# -----------------------------------------------------------------------------
# Extract a frame every SLIDE_INTERVAL seconds (can be float, e.g. 0.5)
# A slide change is detected if SSIM(frame_i, frame_{i-1}) < SIMILARITY.
SLIDE_INTERVAL: float = 1.0     # seconds
SIMILARITY: float = 0.85        # similarity threshold (0–1)
WHISPER_MODEL: str = "base"     # whisper model size (tiny, base, small, medium, large)

# -----------------------------------------------------------------------------
# UTILITY
# -----------------------------------------------------------------------------

def log(message: str) -> None:
    """Timestamped console logger."""
    print(f"[{datetime.now().strftime('%H:%M:%S')}] {message}")

# -----------------------------------------------------------------------------
# AUDIO EXTRACTION
# -----------------------------------------------------------------------------

def extract_audio(video_path: str, output_audio_path: str) -> None:
    """Extract mono 16‑kHz PCM WAV audio from a video file."""
    clip = VideoFileClip(video_path)
    clip.audio.write_audiofile(output_audio_path, fps=16000, codec="pcm_s16le", logger=None)

# -----------------------------------------------------------------------------
# SLIDE DETECTION + EXTRACTION (SINGLE PASS)
# -----------------------------------------------------------------------------

def detect_and_extract_slides(
    video_path: str,
    interval: float,
    threshold: float,
    temp_dir: str,
) -> Tuple[List[float], List[str]]:
    """Single‑pass slide detection and image extraction.

    Returns a tuple (slide_timestamps, slide_image_paths).
    The final timestamp (video end) is appended for alignment.
    """
    log("Detecting and extracting slides")

    clip = VideoFileClip(video_path)
    duration = clip.duration

    # Generate timestamps with float support (e.g. 0.5 s interval)
    timestamps = [round(t, 3) for t in np.arange(0, duration, interval)]

    prev_gray: np.ndarray | None = None
    slide_timestamps: List[float] = []
    slide_images: List[str] = []

    total_frames = len(timestamps)
    for idx, t in enumerate(timestamps):
        frame = clip.get_frame(t)
        gray = cv2.cvtColor(frame, cv2.COLOR_RGB2GRAY)

        is_new_slide = False
        if prev_gray is None:
            is_new_slide = True  # first frame
        else:
            score, _ = ssim(prev_gray, gray, full=True)
            is_new_slide = score < threshold

        if is_new_slide:
            img_path = os.path.join(temp_dir, f"slide_{len(slide_images):02d}.jpg")
            Image.fromarray(frame).save(img_path)
            slide_images.append(img_path)
            slide_timestamps.append(t)

        prev_gray = gray

        # Progress indicator in place
        print(f"\rFrames processed: {idx + 1}/{total_frames} | Slides: {len(slide_images)}", end="", flush=True)

    print()  # newline after progress

    # Ensure final timestamp for transcript alignment
    if not np.isclose(slide_timestamps[-1], duration):
        slide_timestamps.append(round(duration, 3))

    log(f"Slide detection complete: {len(slide_images)} slides located")
    return slide_timestamps, slide_images

# -----------------------------------------------------------------------------
# AUDIO TRANSCRIPTION (WHISPER)
# -----------------------------------------------------------------------------

def transcribe_audio(audio_path: str, model_name: str = WHISPER_MODEL):
    log("Starting Whisper transcription")
    model = whisper.load_model(model_name)
    result = model.transcribe(audio_path, verbose=False)
    log("Transcription complete")
    return [(seg["start"], seg["end"], seg["text"]) for seg in result["segments"]]

# -----------------------------------------------------------------------------
# ALIGN TRANSCRIPT TO SLIDE RANGES
# -----------------------------------------------------------------------------

def group_segments_by_slide(segments, timestamps):
    slide_texts: List[str] = []
    for start, end in zip(timestamps[:-1], timestamps[1:]):
        text = " ".join(seg[2] for seg in segments if start <= seg[0] < end)
        slide_texts.append(text)
    return slide_texts

# -----------------------------------------------------------------------------
# POWERPOINT GENERATION
# -----------------------------------------------------------------------------

def create_ppt(slide_images, slide_notes, output_path):
    log("Creating PowerPoint presentation")
    prs = Presentation()
    blank = prs.slide_layouts[6]

    for img, note in zip(slide_images, slide_notes):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        slide.notes_slide.notes_text_frame.text = note

    prs.save(output_path)
    log("PowerPoint saved")

# -----------------------------------------------------------------------------
# MAIN ENTRY
# -----------------------------------------------------------------------------

def main():
    if len(sys.argv) < 2:
        print("Usage: python vid-to-ppt.py <video_file>")
        sys.exit(1)

    video_path = sys.argv[1]
    if not os.path.exists(video_path):
        print(f"Error: file '{video_path}' not found")
        sys.exit(1)

    base_name = os.path.splitext(os.path.basename(video_path))[0]
    output_pptx = os.path.abspath(f"{base_name}.pptx")

    with tempfile.TemporaryDirectory() as temp_dir:
        audio_path = os.path.join(temp_dir, "audio.wav")

        extract_audio(video_path, audio_path)
        timestamps, slide_images = detect_and_extract_slides(video_path, SLIDE_INTERVAL, SIMILARITY, temp_dir)
        segments = transcribe_audio(audio_path)
        slide_notes = group_segments_by_slide(segments, timestamps)
        create_ppt(slide_images, slide_notes, output_pptx)

    log(f"Process finished. Presentation: {output_pptx}")

if __name__ == "__main__":
    main()
